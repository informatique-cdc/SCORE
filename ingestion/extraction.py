"""
Text extraction from various document formats.

Extracts plain text + heading structure from:
  - HTML (Confluence pages, web content)
  - PDF
  - DOCX / PPTX
  - Markdown
  - Plain text
"""

import io
import logging
import re
from dataclasses import dataclass, field

from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)


@dataclass
class ExtractedText:
    """Result of text extraction with heading structure preserved."""

    text: str
    headings: list[dict] = field(default_factory=list)
    # headings: [{"level": 1, "text": "...", "offset": char_offset}, ...]
    word_count: int = 0
    metadata: dict = field(default_factory=dict)


def extract_text(content: bytes | str, content_type: str) -> ExtractedText:
    """
    Route content to the appropriate extractor based on content type.
    Never log the actual content to avoid data leakage.
    """
    logger.debug("Extracting text from content_type=%s, size=%d", content_type, len(content))

    ct = content_type.lower()
    try:
        if "html" in ct:
            return _extract_html(
                content if isinstance(content, str) else content.decode("utf-8", errors="replace")
            )
        elif "pdf" in ct:
            return _extract_pdf(content if isinstance(content, bytes) else content.encode())
        elif "wordprocessingml" in ct or ct == "application/docx":
            return _extract_docx(content if isinstance(content, bytes) else content.encode())
        elif "presentationml" in ct or ct == "application/pptx":
            return _extract_pptx(content if isinstance(content, bytes) else content.encode())
        elif "markdown" in ct or ct == "text/markdown":
            return _extract_markdown(
                content if isinstance(content, str) else content.decode("utf-8", errors="replace")
            )
        else:
            # Plain text fallback
            text = (
                content if isinstance(content, str) else content.decode("utf-8", errors="replace")
            )
            return ExtractedText(text=text.strip(), word_count=len(text.split()))
    except (ValueError, UnicodeDecodeError, KeyError, OSError) as e:
        logger.error("Text extraction failed for content_type=%s: %s", content_type, e)
        return ExtractedText(text="", metadata={"extraction_error": str(e)})


def _extract_html(html: str) -> ExtractedText:
    """Extract text and heading structure from HTML."""
    soup = BeautifulSoup(html, "html.parser")

    # Remove script/style elements
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()

    headings = []
    text_parts = []
    offset = 0

    for element in soup.descendants:
        if element.name and element.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(element.name[1])
            heading_text = element.get_text(strip=True)
            headings.append({"level": level, "text": heading_text, "offset": offset})
            text_parts.append(f"\n{'#' * level} {heading_text}\n")
            offset += len(text_parts[-1])
        elif element.string and element.parent.name not in ("h1", "h2", "h3", "h4", "h5", "h6"):
            text = element.string.strip()
            if text:
                text_parts.append(text + " ")
                offset += len(text) + 1

    full_text = "".join(text_parts).strip()
    # Collapse excessive whitespace
    full_text = re.sub(r"\n{3,}", "\n\n", full_text)

    return ExtractedText(
        text=full_text,
        headings=headings,
        word_count=len(full_text.split()),
    )


def _extract_pdf(content: bytes) -> ExtractedText:
    """Extract text from PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    text_parts = []
    headings = []
    offset = 0

    for i, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        text_parts.append(page_text)
        offset += len(page_text)

    full_text = "\n\n".join(text_parts).strip()

    # Try to detect headings from outline/bookmarks
    try:
        for item in reader.outline:
            if hasattr(item, "title"):
                headings.append({"level": 1, "text": item.title, "offset": 0})
    except (AttributeError, TypeError, IndexError):
        pass

    return ExtractedText(
        text=full_text,
        headings=headings,
        word_count=len(full_text.split()),
    )


def _extract_docx(content: bytes) -> ExtractedText:
    """Extract text from DOCX with heading structure."""
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(content))
    text_parts = []
    headings = []
    offset = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = (para.style.name or "").lower()
        if "heading" in style_name:
            # Extract heading level from style name (e.g., "Heading 2")
            level_match = re.search(r"(\d+)", style_name)
            level = int(level_match.group(1)) if level_match else 1
            headings.append({"level": level, "text": text, "offset": offset})
            text_parts.append(f"\n{'#' * level} {text}\n")
        else:
            text_parts.append(text + "\n")

        offset += len(text_parts[-1])

    full_text = "".join(text_parts).strip()
    return ExtractedText(text=full_text, headings=headings, word_count=len(full_text.split()))


def _extract_pptx(content: bytes) -> ExtractedText:
    """Extract text from PPTX presentations."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    text_parts = []
    headings = []
    offset = 0

    for i, slide in enumerate(prs.slides, 1):
        slide_header = f"\n## Slide {i}\n"
        headings.append({"level": 2, "text": f"Slide {i}", "offset": offset})
        text_parts.append(slide_header)
        offset += len(slide_header)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        text_parts.append(text + "\n")
                        offset += len(text) + 1

    full_text = "".join(text_parts).strip()
    return ExtractedText(text=full_text, headings=headings, word_count=len(full_text.split()))


def _extract_markdown(md_text: str) -> ExtractedText:
    """Extract headings from Markdown text."""
    headings = []
    offset = 0

    for line in md_text.split("\n"):
        match = re.match(r"^(#{1,6})\s+(.+)$", line)
        if match:
            level = len(match.group(1))
            headings.append({"level": level, "text": match.group(2).strip(), "offset": offset})
        offset += len(line) + 1

    return ExtractedText(
        text=md_text.strip(),
        headings=headings,
        word_count=len(md_text.split()),
    )
